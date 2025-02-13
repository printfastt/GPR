"""
import os
import fnmatch
from pptx import Presentation
from pptx.util import Inches

trial_folder = "depth25to100test"
presentation_date = "02/01/2026"

seq_conf_saves_dir = "C:/Users/carso.LAPTOP/OneDrive/Desktop/Desktop All/research/GPR/SeqConfSaves"
base_dir = os.path.join(seq_conf_saves_dir, trial_folder)
save_dir = os.path.join(os.path.dirname(__file__), "Presentations")
pptx_filename = f"{presentation_date.replace('/', '')}.pptx"
save_path = os.path.join(save_dir, pptx_filename)

image_patterns = [
    'TR_*TE_*.png',
    'TR_*TE_*Original.png',
    'TR_*TE_*ROC.png',
    'TR_*TE_*Thresholded.png'
]

subfolders = [
    os.path.join(base_dir, d)
    for d in os.listdir(base_dir)
    if os.path.isdir(os.path.join(base_dir, d))
]

subfolders.sort(key=lambda x: os.path.getmtime(x))

prs = Presentation()

title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title = title_slide.shapes.title
subtitle = title_slide.placeholders[1]

title.text = f"{presentation_date} Research"
subtitle.text = "Carson Pautz"

overview_slide = prs.slides.add_slide(prs.slide_layouts[5])
text_box = overview_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
text_frame = text_box.text_frame
text_frame.text = "Included Trials:\n"

for subfolder in subfolders:
    text_frame.add_paragraph().text = os.path.basename(subfolder)

for subfolder in subfolders:
    images = {}

    for root, _, files in os.walk(subfolder):
        for pattern in image_patterns:
            match = next((os.path.join(root, f) for f in files if fnmatch.fnmatch(f, pattern)), None)
            if match:
                images[pattern] = match

    if len(images) == 4:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  

        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5), Inches(0.5))
        title_text = title_shape.text_frame
        title_text.text = os.path.basename(subfolder)

        positions = [
            (0.25, 0.82),  # Slightly shifted up
            (5.25, 0.82),  # Slightly shifted up
            (0.25, 4.32),  # Slightly shifted up
            (5.25, 4.32)   # Slightly shifted up
        ]

        for i, pattern in enumerate(image_patterns):
            img_path = images.get(pattern)
            if img_path:
                left, top = positions[i]
                slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(4.5))

if os.path.exists(save_path):
    os.remove(save_path)

prs.save(save_path)
print(f"PowerPoint saved as {save_path}")
"""



import os
import fnmatch
from pptx import Presentation
from pptx.util import Inches

trial_folder1 = "25to100"
trial_folder2 = "0to100"
presentation_date = "02/01/2025"

seq_conf_saves_dir = "C:/Users/carso.LAPTOP/OneDrive/Desktop/Desktop All/research/GPR/SeqConfSaves"
base_dir1 = os.path.join(seq_conf_saves_dir, trial_folder1)
base_dir2 = os.path.join(seq_conf_saves_dir, trial_folder2)

save_dir = os.path.join(os.path.dirname(__file__), "Presentations")
pptx_filename = f"{presentation_date.replace('/', '')}.pptx"
save_path = os.path.join(save_dir, pptx_filename)

image_patterns = [
    'TR_*TE_*.png',
    'TR_*TE_*ROC.png',
    'TR_*TE_*Thresholded.png'
]

subfolders1 = [
    os.path.join(base_dir1, d)
    for d in os.listdir(base_dir1)
    if os.path.isdir(os.path.join(base_dir1, d))
]

subfolders2 = [
    os.path.join(base_dir2, d)
    for d in os.listdir(base_dir2)
    if os.path.isdir(os.path.join(base_dir2, d))
]

subfolders1.sort(key=lambda x: os.path.getmtime(x))
subfolders2.sort(key=lambda x: os.path.getmtime(x))

prs = Presentation()

title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title = title_slide.shapes.title
subtitle = title_slide.placeholders[1]

title.text = f"{presentation_date} Research"
subtitle.text = "Carson Pautz"

for subfolder1, subfolder2 in zip(subfolders1, subfolders2):
    images1 = {}
    images2 = {}

    for root, _, files in os.walk(subfolder1):
        for pattern in image_patterns:
            match = next((os.path.join(root, f) for f in files if fnmatch.fnmatch(f, pattern)), None)
            if match:
                images1[pattern] = match

    for root, _, files in os.walk(subfolder2):
        for pattern in image_patterns:
            match = next((os.path.join(root, f) for f in files if fnmatch.fnmatch(f, pattern)), None)
            if match:
                images2[pattern] = match

    if len(images1) == 3 and len(images2) == 3:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        positions_top = [
            (0.5, 1),
            (3.5, 1),
            (6.5, 1)
        ]

        positions_bottom = [
            (0.5, 4),
            (3.5, 4),
            (6.5, 4)
        ]

        for i, pattern in enumerate(image_patterns):
            img_path1 = images1.get(pattern)
            img_path2 = images2.get(pattern)

            if img_path1:
                left, top = positions_top[i]
                slide.shapes.add_picture(img_path1, Inches(left), Inches(top), width=Inches(2.75))

            if img_path2:
                left, top = positions_bottom[i]
                slide.shapes.add_picture(img_path2, Inches(left), Inches(top), width=Inches(2.75))

if os.path.exists(save_path):
    os.remove(save_path)

prs.save(save_path)
print(f"PowerPoint saved as {save_path}")

